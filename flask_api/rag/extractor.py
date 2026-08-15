"""
extractor.py — Extract text from PPTX files and images.
- PPTX: uses python-pptx to pull all slide text, titles, notes
- Images: uses Groq vision API (llama-4-scout) to describe image content
"""
import os
import base64
import io
from groq import Groq

# ─── PPTX Extraction ─────────────────────────────────────────────────────────

def extract_pptx_text(file_bytes: bytes) -> str:
    """Extract all text content from a PPTX file (bytes)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return ""

    prs = Presentation(io.BytesIO(file_bytes))
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Slide {slide_num} ---"]

        # Extract from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)

            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(f"[Table Row] {row_text}")

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Notes] {notes_text}")

        all_text.extend(slide_texts)

    return "\n".join(all_text)


# ─── Image Extraction ─────────────────────────────────────────────────────────

def extract_image_text(file_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    """
    Use Groq vision model to describe image content.
    Returns a textual description of the image suitable for RAG.
    """
    api_key = os.environ.get("GROQ_API_KEY", "")
    if not api_key:
        return ""

    try:
        client = Groq(api_key=api_key)

        # Encode to base64
        b64_image = base64.b64encode(file_bytes).decode("utf-8")
        data_uri = f"data:{mime_type};base64,{b64_image}"

        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {"url": data_uri}
                        },
                        {
                            "type": "text",
                            "text": (
                                "Describe this image in detail for use in a presentation context. "
                                "Focus on: main subject, key visual elements, text visible in the image, "
                                "color scheme, layout style, design patterns, and any data or charts shown. "
                                "Be thorough and factual — this description will be used to inform AI-generated slides."
                            )
                        }
                    ]
                }
            ],
            max_tokens=800,
            temperature=0.1
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"[RAG] Image extraction failed: {e}")
        return ""
